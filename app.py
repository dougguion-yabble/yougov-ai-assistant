import streamlit as st
import os
import re
import io
import json
import tempfile
from pathlib import Path
from typing import List

# Page config
st.set_page_config(
    page_title="YouGov AI Assistant",
    page_icon="🎯",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Google Drive Folder ID (from your URL)
GDRIVE_FOLDER_ID = "1Ex3CzAA_Xb45CLpO0GRFOLt2cHS9t2l0"

# Supported file extensions
SUPPORTED_EXTENSIONS = {'.txt', '.docx', '.pdf', '.ppt', '.pptx', '.xls', '.xlsx'}

# Custom CSS
st.markdown("""
<style>
    .stApp { background-color: #F8F9FA; }
    .main-header {
        background: linear-gradient(135deg, #2D2D2D, #3D3D3D);
        padding: 1.5rem 2rem;
        border-radius: 12px;
        margin-bottom: 1.5rem;
        color: white;
    }
    .main-header h1 { color: white; margin: 0; font-size: 1.75rem; }
    .main-header .subtitle { color: #E91E8C; font-size: 1rem; margin-top: 0.25rem; }
    .user-message {
        background: linear-gradient(135deg, #E91E8C, #9B4DFF);
        color: white;
        padding: 1rem 1.25rem;
        border-radius: 16px;
        border-bottom-right-radius: 4px;
        margin: 0.5rem 0;
        max-width: 80%;
        margin-left: auto;
    }
    .assistant-message {
        background: white;
        border: 1px solid #E5E7EB;
        padding: 1rem 1.25rem;
        border-radius: 16px;
        border-bottom-left-radius: 4px;
        margin: 0.5rem 0;
        max-width: 85%;
    }
    .sources-box {
        background: rgba(233, 30, 140, 0.05);
        border-left: 3px solid #E91E8C;
        padding: 0.75rem 1rem;
        border-radius: 0 8px 8px 0;
        margin-top: 0.75rem;
        font-size: 0.875rem;
    }
    .sources-label { color: #E91E8C; font-weight: 600; font-size: 0.75rem; margin-bottom: 0.25rem; }
    .sidebar-stat {
        background: white;
        padding: 1rem;
        border-radius: 8px;
        text-align: center;
        border: 1px solid #E5E7EB;
        margin-bottom: 0.5rem;
    }
    .stat-value { font-size: 1.5rem; font-weight: 700; color: #E91E8C; }
    .stat-label { font-size: 0.75rem; color: #6B7280; }
    #MainMenu {visibility: hidden;}
    footer {visibility: hidden;}
</style>
""", unsafe_allow_html=True)

# Initialize session state
if 'messages' not in st.session_state:
    st.session_state.messages = []
if 'knowledge_base_ready' not in st.session_state:
    st.session_state.knowledge_base_ready = False
if 'collection' not in st.session_state:
    st.session_state.collection = None
if 'embedding_model' not in st.session_state:
    st.session_state.embedding_model = None
if 'gdrive_loaded' not in st.session_state:
    st.session_state.gdrive_loaded = False


@st.cache_resource
def load_embedding_model():
    from sentence_transformers import SentenceTransformer
    return SentenceTransformer('all-MiniLM-L6-v2')


@st.cache_resource
def init_chromadb():
    import chromadb
    from chromadb.config import Settings
    client = chromadb.Client(Settings(anonymized_telemetry=False))
    collection = client.get_or_create_collection(name="yougov_knowledge")
    return client, collection


def get_gdrive_service():
    """Initialize Google Drive service using credentials from Streamlit secrets."""
    from google.oauth2 import service_account
    from googleapiclient.discovery import build
    
    # Get credentials from Streamlit secrets
    credentials_dict = dict(st.secrets["gcp_service_account"])
    credentials = service_account.Credentials.from_service_account_info(
        credentials_dict,
        scopes=['https://www.googleapis.com/auth/drive.readonly']
    )
    service = build('drive', 'v3', credentials=credentials)
    return service


def list_files_in_folder(service, folder_id):
    """List all supported files in the Google Drive folder."""
    query = f"'{folder_id}' in parents and trashed=false"
    results = service.files().list(
        q=query,
        fields="files(id, name, mimeType)",
        pageSize=100
    ).execute()
    
    files = results.get('files', [])
    # Filter for supported file types
    supported_files = []
    for f in files:
        ext = Path(f['name']).suffix.lower()
        if ext in SUPPORTED_EXTENSIONS:
            supported_files.append(f)
    
    return supported_files


def download_file(service, file_id, filename):
    """Download a file from Google Drive."""
    from googleapiclient.http import MediaIoBaseDownload
    
    request = service.files().get_media(fileId=file_id)
    content = io.BytesIO()
    downloader = MediaIoBaseDownload(content, request)
    
    done = False
    while not done:
        status, done = downloader.next_chunk()
    
    content.seek(0)
    return content.read()


def extract_text_from_txt(file_content: bytes, filename: str) -> str:
    return file_content.decode('utf-8', errors='ignore')


def extract_text_from_docx(file_content: bytes, filename: str) -> str:
    from docx import Document
    doc = Document(io.BytesIO(file_content))
    return '\n\n'.join([para.text for para in doc.paragraphs if para.text.strip()])


def extract_text_from_pdf(file_content: bytes, filename: str) -> str:
    import pdfplumber
    text_parts = []
    with pdfplumber.open(io.BytesIO(file_content)) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                text_parts.append(page_text)
    return '\n\n'.join(text_parts)


def extract_text_from_pptx(file_content: bytes, filename: str) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(file_content))
    text_parts = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_text = [f"[Slide {slide_num}]"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text)
        if len(slide_text) > 1:
            text_parts.append('\n'.join(slide_text))
    return '\n\n'.join(text_parts)


def extract_text_from_excel(file_content: bytes, filename: str) -> str:
    ext = Path(filename).suffix.lower()
    text_parts = []
    if ext == '.xlsx':
        from openpyxl import load_workbook
        wb = load_workbook(io.BytesIO(file_content), data_only=True)
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            sheet_text = [f"[Sheet: {sheet_name}]"]
            for row in sheet.iter_rows():
                row_values = [str(cell.value) for cell in row if cell.value is not None]
                if row_values:
                    sheet_text.append(' | '.join(row_values))
            if len(sheet_text) > 1:
                text_parts.append('\n'.join(sheet_text))
        wb.close()
    elif ext == '.xls':
        import xlrd
        wb = xlrd.open_workbook(file_contents=file_content)
        for sheet_idx in range(wb.nsheets):
            sheet = wb.sheet_by_index(sheet_idx)
            sheet_text = [f"[Sheet: {sheet.name}]"]
            for row_idx in range(sheet.nrows):
                row_values = [str(sheet.cell_value(row_idx, col_idx)) for col_idx in range(sheet.ncols) if sheet.cell_value(row_idx, col_idx)]
                if row_values:
                    sheet_text.append(' | '.join(row_values))
            if len(sheet_text) > 1:
                text_parts.append('\n'.join(sheet_text))
    return '\n\n'.join(text_parts)


def extract_text(file_content: bytes, filename: str) -> str:
    ext = Path(filename).suffix.lower()
    extractors = {
        '.txt': extract_text_from_txt,
        '.docx': extract_text_from_docx,
        '.pdf': extract_text_from_pdf,
        '.ppt': extract_text_from_pptx,
        '.pptx': extract_text_from_pptx,
        '.xls': extract_text_from_excel,
        '.xlsx': extract_text_from_excel,
    }
    if ext not in extractors:
        raise ValueError(f"Unsupported file format: {ext}")
    return extractors[ext](file_content, filename)


def chunk_text(text: str, chunk_size: int = 500) -> List[str]:
    sentences = re.split(r'(?<=[.!?])\s+', text)
    chunks = []
    current_chunk = []
    current_word_count = 0
    for sentence in sentences:
        sentence_words = len(sentence.split())
        if current_word_count + sentence_words > chunk_size and current_chunk:
            chunks.append(' '.join(current_chunk))
            current_chunk = current_chunk[-2:] if len(current_chunk) > 2 else []
            current_word_count = sum(len(s.split()) for s in current_chunk)
        current_chunk.append(sentence)
        current_word_count += sentence_words
    if current_chunk:
        chunks.append(' '.join(current_chunk))
    return chunks


def generate_answer(question: str, context_chunks: list, sources: list, api_key: str) -> str:
    import anthropic
    client = anthropic.Anthropic(api_key=api_key)
    
    context_text = ""
    for i, (chunk, source) in enumerate(zip(context_chunks, sources), 1):
        context_text += f"\n[Source {i}: {source}]\n{chunk}\n"

    system_prompt = """You are a helpful assistant for YouGov/Yabble. Answer questions based ONLY on the provided context.

RULES:
1. ONLY use information from the provided context - never make things up
2. Write a clear, professional, well-structured response
3. Use bullet points where appropriate for readability
4. If the context doesn't fully answer the question, say what you can answer and note what's missing
5. Write in a tone suitable for sales/account management professionals
6. Keep the response concise but complete - aim for 2-4 paragraphs
7. Do NOT include source citations inline - they will be added separately at the end"""

    user_prompt = f"""Based on the following context from our knowledge base, please answer this question:

QUESTION: {question}

CONTEXT:
{context_text}

Please provide a clear, well-structured answer based only on the context above."""

    response = client.messages.create(
        model="claude-sonnet-4-20250514",
        max_tokens=1024,
        messages=[{"role": "user", "content": user_prompt}],
        system=system_prompt
    )
    return response.content[0].text


def load_documents_from_gdrive(embedding_model, chroma_client):
    """Load all documents from Google Drive folder."""
    try:
        service = get_gdrive_service()
        files = list_files_in_folder(service, GDRIVE_FOLDER_ID)
        
        if not files:
            return None, 0, []
        
        # Clear existing collection
        try:
            chroma_client.delete_collection("yougov_knowledge")
        except:
            pass
        collection = chroma_client.get_or_create_collection(name="yougov_knowledge")
        
        total_chunks = 0
        processed_files = []
        progress_bar = st.progress(0)
        status_text = st.empty()
        
        for idx, file_info in enumerate(files):
            filename = file_info['name']
            file_id = file_info['id']
            
            status_text.text(f"Processing: {filename}")
            
            try:
                # Download file
                file_content = download_file(service, file_id, filename)
                
                # Extract and chunk
                text = extract_text(file_content, filename)
                chunks = chunk_text(text)
                
                if chunks:
                    ids = [f"{filename}_{i}" for i in range(len(chunks))]
                    metadatas = [{"source": filename, "chunk_index": i} for i in range(len(chunks))]
                    embeddings = embedding_model.encode(chunks).tolist()
                    
                    collection.add(
                        documents=chunks,
                        embeddings=embeddings,
                        metadatas=metadatas,
                        ids=ids
                    )
                    total_chunks += len(chunks)
                    processed_files.append(filename)
                    
            except Exception as e:
                st.warning(f"Error processing {filename}: {e}")
            
            progress_bar.progress((idx + 1) / len(files))
        
        status_text.empty()
        progress_bar.empty()
        
        return collection, total_chunks, processed_files
        
    except Exception as e:
        st.error(f"Error connecting to Google Drive: {e}")
        return None, 0, []


# Header
st.markdown("""
<div class="main-header">
    <h1>🎯 YouGov AI Assistant</h1>
    <div class="subtitle">Ask questions about Yabble & AI capabilities</div>
</div>
""", unsafe_allow_html=True)

# Sidebar
with st.sidebar:
    st.markdown("### ⚙️ Setup")
    
    api_key = st.text_input("Anthropic API Key", type="password", help="Enter your Anthropic API key")
    
    if api_key:
        st.success("✅ API Key set")
    
    st.markdown("---")
    st.markdown("### 📁 Knowledge Base")
    
    # Load from Google Drive button
    if st.button("🔄 Load Documents from Google Drive", type="primary"):
        with st.spinner("Loading documents from Google Drive..."):
            embedding_model = load_embedding_model()
            chroma_client, _ = init_chromadb()
            
            collection, total_chunks, processed_files = load_documents_from_gdrive(
                embedding_model, chroma_client
            )
            
            if collection and total_chunks > 0:
                st.session_state.collection = collection
                st.session_state.embedding_model = embedding_model
                st.session_state.knowledge_base_ready = True
                st.session_state.gdrive_loaded = True
                st.success(f"✅ Loaded {len(processed_files)} files, {total_chunks} chunks")
            else:
                st.error("No documents found or error loading")
    
    st.markdown("---")
    st.markdown("### 📊 Status")
    
    if st.session_state.knowledge_base_ready and st.session_state.collection:
        count = st.session_state.collection.count()
        st.markdown(f"""
        <div class="sidebar-stat">
            <div class="stat-value">{count}</div>
            <div class="stat-label">Chunks Indexed</div>
        </div>
        """, unsafe_allow_html=True)
        
        if count > 0:
            all_data = st.session_state.collection.get(include=['metadatas'])
            sources = set(m['source'] for m in all_data['metadatas'])
            st.markdown(f"**{len(sources)} documents loaded**")
            
            with st.expander("View documents"):
                for source in sorted(sources):
                    st.markdown(f"- {source}")
    else:
        st.info("Click 'Load Documents' to get started")


# Main chat area
if not api_key:
    st.warning("👈 Enter your Anthropic API Key in the sidebar")
    st.markdown("""
    **Need an API key?**
    1. Go to [console.anthropic.com](https://console.anthropic.com)
    2. Create an account
    3. Generate an API key
    4. Paste it in the sidebar
    """)
elif not st.session_state.knowledge_base_ready:
    st.info("👈 Click **'Load Documents from Google Drive'** in the sidebar to get started.")
else:
    # Display chat history
    for message in st.session_state.messages:
        if message["role"] == "user":
            st.markdown(f'<div class="user-message">{message["content"]}</div>', unsafe_allow_html=True)
        else:
            st.markdown(f'<div class="assistant-message">{message["content"]}</div>', unsafe_allow_html=True)
            if message.get("sources"):
                sources_html = "<br>".join([f"• {s}" for s in message["sources"]])
                st.markdown(f'''
                <div class="sources-box">
                    <div class="sources-label">📎 SOURCES</div>
                    {sources_html}
                </div>
                ''', unsafe_allow_html=True)
    
    # Chat input
    if prompt := st.chat_input("Ask a question about Yabble or AI..."):
        st.session_state.messages.append({"role": "user", "content": prompt})
        st.markdown(f'<div class="user-message">{prompt}</div>', unsafe_allow_html=True)
        
        with st.spinner("Generating answer..."):
            query_embedding = st.session_state.embedding_model.encode([prompt]).tolist()
            results = st.session_state.collection.query(
                query_embeddings=query_embedding,
                n_results=5,
                include=['documents', 'metadatas', 'distances']
            )
            
            if results['documents'] and results['documents'][0]:
                chunks = results['documents'][0]
                sources = [m['source'] for m in results['metadatas'][0]]
                unique_sources = list(dict.fromkeys(sources))
                
                answer = generate_answer(prompt, chunks, sources, api_key)
                
                st.markdown(f'<div class="assistant-message">{answer}</div>', unsafe_allow_html=True)
                sources_html = "<br>".join([f"• {s}" for s in unique_sources])
                st.markdown(f'''
                <div class="sources-box">
                    <div class="sources-label">📎 SOURCES</div>
                    {sources_html}
                </div>
                ''', unsafe_allow_html=True)
                
                st.session_state.messages.append({
                    "role": "assistant",
                    "content": answer,
                    "sources": unique_sources
                })
            else:
                no_answer = "I couldn't find any relevant information in the knowledge base for this question."
                st.markdown(f'<div class="assistant-message">{no_answer}</div>', unsafe_allow_html=True)
                st.session_state.messages.append({"role": "assistant", "content": no_answer})
